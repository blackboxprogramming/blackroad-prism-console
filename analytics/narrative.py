import json
from datetime import datetime
from pathlib import Path
from typing import List

from .utils import increment, log_event, validate

ROOT = Path(__file__).resolve().parents[1]
ART = ROOT / "artifacts" / "reports"


def build_report(plan_path: Path, out_prefix: Path) -> None:
    plan = json.loads(plan_path.read_text())
    anomalies_path = ROOT / "artifacts" / "anomalies" / "latest.json"
    cohorts_dir = ROOT / "artifacts" / "cohorts"
    anomalies = json.loads(anomalies_path.read_text()) if anomalies_path.exists() else []
    cohort_files = list(cohorts_dir.glob("*.json"))
    cohort_names = [f.stem for f in cohort_files]
    sections = [
        {"title": "What happened", "body": f"Anomalies: {len(anomalies)}"},
        {"title": "Why it matters", "body": "Impacting key KPIs"},
        {"title": "What we're doing", "body": ", ".join(a["action"] for a in plan.get("actions", [])) or "No actions"},
        {"title": "Risks & Next Steps", "body": "Monitor cohorts: " + ", ".join(cohort_names)},
    ]
    data = {"sections": sections}
    validate(data, "narrative.schema.json")
    ts = datetime.utcnow().strftime("%Y%m%d%H%M%S")
    ART.mkdir(parents=True, exist_ok=True)
    out_prefix.parent.mkdir(parents=True, exist_ok=True)
    md_path = out_prefix.with_suffix(".md")
    md_lines = [f"# Executive Report {ts}", ""]
    for s in sections:
        md_lines.append(f"## {s['title']}")
        md_lines.append(s["body"])
        md_lines.append("")
    md_path.write_text("\n".join(md_lines))
    try:
        from pptx import Presentation  # type: ignore

        ppt_path = out_prefix.with_suffix(".pptx")
        pres = Presentation()
        for s in sections:
            slide = pres.slides.add_slide(pres.slide_layouts[1])
            slide.shapes.title.text = s["title"]
            slide.shapes.placeholders[1].text = s["body"]
        pres.save(ppt_path)
    except Exception:
        slides_md = out_prefix.with_name(out_prefix.name + "_slides.md")
        slide_lines: List[str] = []
        for s in sections:
            slide_lines.append(f"# {s['title']}")
            slide_lines.append(s["body"])
            slide_lines.append("")
        slides_md.write_text("\n".join(slide_lines))
    (out_prefix.with_suffix(".json")).write_text(json.dumps(data, indent=2))
    increment("narrative_built")
    log_event({"type": "narrative_built", "plan": str(plan_path)})
