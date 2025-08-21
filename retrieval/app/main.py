import io
import os
import hashlib
from typing import List, Dict, Optional, Any
from fastapi import FastAPI, UploadFile, File, Depends, HTTPException, Security
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPAuthorizationCredentials, HTTPBearer
from pydantic import BaseModel, Field
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.http.models import (
    Distance, VectorParams, PointStruct, Filter, FieldCondition, MatchValue
)
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from bs4 import BeautifulSoup

# --------------------------
# Env & Config
# --------------------------
BEARER_TOKEN = os.getenv("BEARER_TOKEN", "change-me")
QDRANT_URL = os.getenv("QDRANT_URL", "http://localhost:6333")
QDRANT_API_KEY = os.getenv("QDRANT_API_KEY", None) or None
QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION", "blackroad_docs")

EMBEDDING_MODEL_NAME = os.getenv("EMBEDDING_MODEL", "BAAI/bge-small-en-v1.5")
EMBEDDING_DEVICE = os.getenv("EMBEDDING_DEVICE", "cpu")

CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "1000"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "200"))
MAX_QUERY_TOP_K = int(os.getenv("MAX_QUERY_TOP_K", "12"))

ALLOW_ORIGINS = os.getenv("CORS_ALLOW_ORIGINS", "*")

# --------------------------
# Auth
# --------------------------
bearer_scheme = HTTPBearer(auto_error=False)
def verify_auth(creds: Optional[HTTPAuthorizationCredentials] = Security(bearer_scheme)):
    if not BEARER_TOKEN:
        raise HTTPException(status_code=500, detail="Server auth misconfigured.")
    if creds is None or creds.scheme.lower() != "bearer" or creds.credentials != BEARER_TOKEN:
        raise HTTPException(status_code=401, detail="Unauthorized")
    return True

# --------------------------
# Data Models
# --------------------------
class Document(BaseModel):
    id: Optional[str] = None
    text: str
    metadata: Dict[str, Any] = Field(default_factory=dict)

class UpsertRequest(BaseModel):
    documents: List[Document]

class QueryItem(BaseModel):
    query: str
    top_k: int = 5
    filter: Optional[Dict[str, Any]] = None

class QueryRequest(BaseModel):
    queries: List[QueryItem]

class DeleteRequest(BaseModel):
    ids: Optional[List[str]] = None
    filter: Optional[Dict[str, Any]] = None
    delete_all: Optional[bool] = False

class ScoredChunk(BaseModel):
    id: str
    score: float
    text: str
    metadata: Dict[str, Any] = Field(default_factory=dict)

class QueryResponseItem(BaseModel):
    query: str
    results: List[ScoredChunk]

class QueryResponse(BaseModel):
    results: List[QueryResponseItem]

# --------------------------
# Utilities
# --------------------------
def sha1(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8")).hexdigest()

def clean_text(s: str) -> str:
    return "\n".join([line.strip() for line in s.replace("\r","").split("\n") if line.strip()])

def chunk_text(s: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    s = s.strip()
    if len(s) <= size:
        return [s] if s else []
    chunks = []
    start = 0
    while start < len(s):
        end = min(start + size, len(s))
        chunks.append(s[start:end])
        if end == len(s): break
        start = max(0, end - overlap)
    return chunks

# --------------------------
# Embeddings
# --------------------------
class Embeddings:
    _model: SentenceTransformer = None
    _dim: int = None

    @classmethod
    def load(cls):
        if cls._model is None:
            cls._model = SentenceTransformer(EMBEDDING_MODEL_NAME, device=EMBEDDING_DEVICE)
            cls._dim = cls._model.get_sentence_embedding_dimension()

    @classmethod
    def encode_docs(cls, texts: List[str]) -> List[List[float]]:
        cls.load()
        # BGE family prefers normalized embeddings
        return cls._model.encode(texts, normalize_embeddings=True).tolist()

    @classmethod
    def encode_query(cls, q: str) -> List[float]:
        cls.load()
        # BGE instructs to prefix queries with "query: "
        return cls._model.encode([f"query: {q}"], normalize_embeddings=True)[0].tolist()

    @classmethod
    def dim(cls) -> int:
        cls.load()
        return cls._dim

# --------------------------
# Vector Store (Qdrant)
# --------------------------
class QdrantStore:
    def __init__(self):
        self.client = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)
        self.ensure_collection()

    def ensure_collection(self):
        existing = [c.name for c in self.client.get_collections().collections]
        if QDRANT_COLLECTION not in existing:
            self.client.create_collection(
                collection_name=QDRANT_COLLECTION,
                vectors_config=VectorParams(size=Embeddings.dim(), distance=Distance.COSINE),
            )

    def _to_filter(self, filt: Optional[Dict[str, Any]]) -> Optional[Filter]:
        if not filt: return None
        must = []
        for k, v in filt.items():
            must.append(FieldCondition(key=f"metadata.{k}", match=MatchValue(value=v)))
        return Filter(must=must) if must else None

    def upsert_chunks(self, items: List[Dict[str, Any]]):
        points = []
        for it in items:
            points.append(PointStruct(
                id=it["id"],
                vector=it["vector"],
                payload={
                    "text": it["text"],
                    "document_id": it.get("document_id"),
                    "metadata": it.get("metadata", {})
                }
            ))
        self.client.upsert(collection_name=QDRANT_COLLECTION, points=points)

    def query(self, query_vec: List[float], top_k: int, filt: Optional[Dict[str, Any]]):
        qf = self._to_filter(filt)
        return self.client.search(
            collection_name=QDRANT_COLLECTION,
            query_vector=query_vec,
            limit=top_k,
            query_filter=qf
        )

    def delete(self, ids: Optional[List[str]] = None, filt: Optional[Dict[str, Any]] = None, delete_all: bool = False):
        if delete_all:
            self.client.delete_collection(QDRANT_COLLECTION)
            self.ensure_collection()
            return
        if ids:
            self.client.delete(collection_name=QDRANT_COLLECTION, points_selector=ids)
        if filt:
            qf = self._to_filter(filt)
            self.client.delete(collection_name=QDRANT_COLLECTION, points_selector=qf)

store: Optional[QdrantStore] = None

# --------------------------
# File Parsers
# --------------------------
def read_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def read_docx(file_bytes: bytes) -> str:
    doc = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs)

def read_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def read_html(file_bytes: bytes) -> str:
    soup = BeautifulSoup(file_bytes.decode("utf-8", errors="ignore"), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n")

def sniff_and_read(filename: str, file_bytes: bytes) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):   return read_pdf(file_bytes)
    if name.endswith(".docx"):  return read_docx(file_bytes)
    if name.endswith(".pptx"):  return read_pptx(file_bytes)
    if name.endswith(".html") or name.endswith(".htm"): return read_html(file_bytes)
    # default: treat as text/markdown
    return file_bytes.decode("utf-8", errors="ignore")

# --------------------------
# FastAPI app
# --------------------------
app = FastAPI(title="BlackRoad × Lucidia Retrieval", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"] if ALLOW_ORIGINS == "*" else [o.strip() for o in ALLOW_ORIGINS.split(",")],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("startup")
def _startup():
    global store
    Embeddings.load()
    store = QdrantStore()

@app.get("/healthz")
def healthz():
    return {"status": "ok", "collection": QDRANT_COLLECTION, "embedding_dim": Embeddings.dim()}

@app.post("/upsert", dependencies=[Depends(verify_auth)])
def upsert(req: UpsertRequest):
    items = []
    texts = []
    idxs = []
    for d in req.documents:
        doc_id = d.id or sha1(d.text)[:16]
        for i, ch in enumerate(chunk_text(clean_text(d.text))):
            chunk_id = f"{doc_id}:{i:04d}"
            items.append({"id": chunk_id, "text": ch, "document_id": doc_id, "metadata": d.metadata})
            texts.append(ch)
            idxs.append(chunk_id)
    if not items:
        return {"upserted": 0}
    vecs = Embeddings.encode_docs(texts)
    for i, v in enumerate(vecs):
        items[i]["vector"] = v
    store.upsert_chunks(items)
    return {"upserted": len(items)}

@app.post("/upsert-file", dependencies=[Depends(verify_auth)])
async def upsert_file(files: List[UploadFile] = File(...)):
    total = 0
    for f in files:
        content = await f.read()
        text = clean_text(sniff_and_read(f.filename, content))
        doc_id = sha1(f.filename + str(len(content)))[:16]
        chunks = chunk_text(text)
        vecs = Embeddings.encode_docs(chunks) if chunks else []
        to_upsert = []
        for i, (ch, vec) in enumerate(zip(chunks, vecs)):
            chunk_id = f"{doc_id}:{i:04d}"
            to_upsert.append({
                "id": chunk_id,
                "text": ch,
                "vector": vec,
                "document_id": doc_id,
                "metadata": {"source": "file", "filename": f.filename}
            })
        if to_upsert:
            store.upsert_chunks(to_upsert)
            total += len(to_upsert)
    return {"upserted": total}

@app.post("/query", response_model=QueryResponse, dependencies=[Depends(verify_auth)])
def query(req: QueryRequest):
    out: List[QueryResponseItem] = []
    for q in req.queries:
        top_k = min(max(q.top_k, 1), MAX_QUERY_TOP_K)
        qvec = Embeddings.encode_query(q.query)
        matches = store.query(qvec, top_k=top_k, filt=q.filter)
        results = []
        for m in matches:
            payload = m.payload or {}
            results.append(ScoredChunk(
                id=str(m.id),
                score=float(m.score),
                text=str(payload.get("text","")),
                metadata=payload.get("metadata", {}) | {"document_id": payload.get("document_id")}
            ))
        out.append(QueryResponseItem(query=q.query, results=results))
    return QueryResponse(results=out)

@app.post("/delete", dependencies=[Depends(verify_auth)])
def delete(req: DeleteRequest):
    store.delete(ids=req.ids, filt=req.filter, delete_all=bool(req.delete_all))
    return {"status": "ok"}
